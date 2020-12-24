import os
from . import helper_functions
from pptx import Presentation
# python-pptx documentation: https://python-pptx.readthedocs.io/en/latest/index.html

# TODO: Clean up this script a bit - There are probably some better ways to handle directory navigation and file saving. For now it works, but it could definitely be better.

# Open the PowerPoint Template
# TODO: Get rid of global variables and make this a function - there is some variable shadowing occurring because of how terrible this is written O.o
templateName = "Risk-Map-Slides-Template.pptx"
cwd = os.getcwd()
prsPath = os.path.join(cwd, 'helper_files', templateName)
prs = Presentation(prsPath)


class CreatePresentation:
    # TODO: This doesn't need to be a class - convert to a function
    def __init__(self, data):
        # data[0] = Slide Layout
        # Set slide layout and create slide
        self.slide = prs.slides.add_slide(prs.slide_layouts[data[0]])

        if data[0] == 0:
            title = self.slide.shapes.title
            subtitle = self.slide.placeholders[1]
            title.text = data[1]
            subtitle.text = data[2]
        elif data[0] == 8:
            title = self.slide.shapes.title
            title.text = data[1]
            imagePlaceholder = self.slide.placeholders[1]
            imagePlaceholder.crop_top = -1
            imagePlaceholder.crop_bottom = -1
            imagePlaceholder.insert_picture(data[3])


def buildSlides():
    prsSlides = [[0, "JTF-COVID J2", "Regional ICU Risk Analysis", ""]]
    cwd = os.getcwd()
    todayImagesPath = os.path.join(cwd, "Risk_Map_Images", helper_functions.get_Date())
    for imageFile in os.listdir(todayImagesPath):
        if imageFile == f"Capital-Region_Risk-Map_{helper_functions.get_Date()}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Capital Region", "", imageFilePath])
        elif imageFile == f"Central-New-York_Risk-Map_{helper_functions.get_Date()}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Central New York", "", imageFilePath])
        elif imageFile == f"Finger-Lakes_Risk-Map_{helper_functions.get_Date()}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Finger Lakes", "", imageFilePath])
        elif imageFile == f"Long-Island_Risk-Map_{helper_functions.get_Date()}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Long Island", "", imageFilePath])
        elif imageFile == f"Mid-Hudson_Risk-Map_{helper_functions.get_Date()}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Mid-Hudson", "", imageFilePath])
        elif imageFile == f"Mohawk-Valley_Risk-Map_{helper_functions.get_Date()}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Mohawk Valley", "", imageFilePath])
        elif imageFile == f"North-Country_Risk-Map_{helper_functions.get_Date()}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "North Country", "", imageFilePath])
        elif imageFile == f"Southern-Tier_Risk-Map_{helper_functions.get_Date()}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Southern Tier", "", imageFilePath])
        elif imageFile == f"Western-New-York_Risk-Map_{helper_functions.get_Date()}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Western New York", "", imageFilePath])
    return prsSlides


def savePresentation(slides):

    """ Slide Layout Reference:
    0 ->  title and subtitle
    1 ->  title and content
    2 ->  section header
    3 ->  two content
    4 ->  Comparison
    5 ->  Title only
    6 ->  Blank
    7 ->  Content with caption
    8 ->  Pic with caption
    """

    """ Class Data Format
    0 -> Slide Layout
    1 -> Slide Title
    2 -> Slide Subtitle
    3 -> Slide Image Path
    3 -> 
    """

    for each_slide in slides:
        CreatePresentation(each_slide)

    # Save PowerPoint
    pptxFileName = f"Regional-Risk-Analysis_{helper_functions.get_Date()}.pptx"
    pptxFilePath = os.path.join(cwd, "../Daily Slides", pptxFileName)
    prs.save(pptxFilePath)

    # Open PowerPoint when finished
    # os.startfile(pptxFilePath)


def makePresentation():
    newSlides = buildSlides()
    savePresentation(newSlides)
