import os
import pytz
from datetime import datetime
from pptx import Presentation


# python-pptx documentation: https://python-pptx.readthedocs.io/en/latest/index.html

"""
This snippet uses the pytz python library to easily ensure we are using EST time
(probs could do this with datetime.timedelta). This is likely only needed if you are working on this notebook
while it is located on a remote server (such as Google Colab)
"""
utc = pytz.utc.localize(datetime.utcnow())
est = utc.astimezone(pytz.timezone("America/New_York"))
date = est.strftime("%m-%d-%Y")
figureDate = est.strftime("%m-%d-%Y %H:%M")

# Create some global variables
cwd = os.getcwd()

# Open the PowerPoint Template
templateName = "Risk-Map-Slides-Template.pptx"
prsPath = os.path.join(cwd, templateName)
prs = Presentation(prsPath)


class CreatePresentation:
    def __init__(self, data):
        # data[0] = Slide Layout
        # Set slide layout and create slide
        self.slide = prs.slides.add_slide(prs.slide_layouts[data[0]])

        if data[0] == 0:
            title = self.slide.shapes.title
            subtitle = self.slide.placeholders[1]
            title.text = data[1]
            subtitle.text = data[2]
        elif data[0] == 8:
            title = self.slide.shapes.title
            title.text = data[1]
            imagePlaceholder = self.slide.placeholders[1]
            imagePlaceholder.crop_top = -1
            imagePlaceholder.crop_bottom = -1
            imagePlaceholder.insert_picture(data[3])


def buildSlides():
    prsSlides = [[0, "JTF-COVID J2", "Regional ICU Risk Analysis", ""]]
    todayImagesPath = os.path.join(cwd, "../Risk_Map_Images", date)
    for imageFile in os.listdir(todayImagesPath):
        if imageFile == f"Capital-Region_Risk-Map_{date}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Capital Region", "", imageFilePath])
        elif imageFile == f"Central-New-York_Risk-Map_{date}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Central New York", "", imageFilePath])
        elif imageFile == f"Finger-Lakes_Risk-Map_{date}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Finger Lakes", "", imageFilePath])
        elif imageFile == f"Long-Island_Risk-Map_{date}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Long Island", "", imageFilePath])
        elif imageFile == f"Mid-Hudson_Risk-Map_{date}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Mid-Hudson", "", imageFilePath])
        elif imageFile == f"Mohawk-Valley_Risk-Map_{date}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "Mohawk Valley", "", imageFilePath])
        elif imageFile == f"North-Country_Risk-Map_{date}.png":
            imageFilePath = os.path.join(todayImagesPath, imageFile)
            prsSlides.append([8, "North Country", "", imageFilePath])
    return prsSlides


def makePresentation(slides):

    """ Slide Layout Reference:  
    0 ->  title and subtitle 
    1 ->  title and content 
    2 ->  section header 
    3 ->  two content 
    4 ->  Comparison 
    5 ->  Title only  
    6 ->  Blank
    7 ->  Content with caption 
    8 ->  Pic with caption 
    """

    """ Class Data Format
    0 -> Slide Layout
    1 -> Slide Title
    2 -> Slide Subtitle
    3 -> Slide Image Path
    3 -> 
    """

    for each_slide in slides:
        CreatePresentation(each_slide)

    # Save PowerPoint
    pptxFileName = f"Regional-Risk-Analysis_{date}.pptx"
    pptxFilePath = os.path.join(cwd, "../Daily Slides", pptxFileName)
    prs.save(pptxFilePath)

    # Open PowerPoint when finished
    os.startfile(pptxFilePath)


newSlides = buildSlides()

makePresentation(newSlides)
